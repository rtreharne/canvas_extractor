import os
import requests
import json
import mimetypes
from io import BytesIO
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
from dotenv import load_dotenv

# --- LOAD ENVIRONMENT VARIABLES ---
load_dotenv()
API_BASE_URL = os.getenv("CANVAS_API_URL") + "/api/v1"
API_TOKEN = os.getenv("CANVAS_API_TOKEN")

try:
    COURSE_ID = int(input("Enter Canvas course ID: ").strip())
except ValueError:
    print("❌ Invalid course ID. Please enter a number.")
    exit(1)

OUTPUT_DIR = "extracted_files"

os.makedirs(OUTPUT_DIR, exist_ok=True)

HEADERS = {
    "Authorization": f"Bearer {API_TOKEN}"
}

# --- UTILITY FUNCTIONS ---

def list_course_files(course_id):
    files = []
    url = f"{API_BASE_URL}/courses/{course_id}/files"
    params = {"per_page": 100}
    while url:
        res = requests.get(url, headers=HEADERS, params=params)
        res.raise_for_status()
        files.extend(res.json())
        url = res.links.get("next", {}).get("url")
    return files

def download_file(file_obj):
    download_url = file_obj["url"]
    response = requests.get(download_url, headers=HEADERS)
    response.raise_for_status()
    return BytesIO(response.content)

def extract_text_from_pdf(file_stream):
    text = ""
    with fitz.open(stream=file_stream, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(file_stream):
    doc = Document(file_stream)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_stream):
    prs = Presentation(file_stream)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def process_file(file_obj):
    filename = file_obj["display_name"]
    ext = os.path.splitext(filename)[1].lower()
    try:
        file_stream = download_file(file_obj)

        if ext == ".pdf":
            text = extract_text_from_pdf(file_stream)
        elif ext == ".docx":
            text = extract_text_from_docx(file_stream)
        elif ext == ".pptx":
            text = extract_text_from_pptx(file_stream)
        else:
            return  # unsupported type

        output_path = os.path.join(OUTPUT_DIR, f"{file_obj['id']}.json")
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump({
                "filename": filename,
                "canvas_file_id": file_obj["id"],
                "text": text.strip()
            }, f, indent=2)

        print(f"✔ Saved: {output_path}")

    except Exception as e:
        print(f"✖ Failed to process {filename}: {e}")

# --- MAIN EXECUTION ---

def main():
    files = list_course_files(COURSE_ID)
    target_exts = [".pdf", ".docx", ".pptx"]
    filtered = [f for f in files if os.path.splitext(f["display_name"])[1].lower() in target_exts]

    print(f"Found {len(filtered)} target files.")
    for f in filtered:
        process_file(f)

if __name__ == "__main__":
    main()
